"""
Generate a professional PPT for Claude Code installation tutorial.
Includes embedded images with proper scaling and clean layout.
Optimized: no overflow, balanced spacing, harmonious colors, proper image sizing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ── Paths ──────────────────────────────────────────────────────
REPO_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(REPO_DIR, "images")
OUTPUT = os.path.join(REPO_DIR, "Claude_Code_安装教程.pptx")

# ── Slide dimensions ───────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.80)
MARGIN_R = Inches(0.80)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R         # ≈ 11.73"

# ── Color palette ──────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x14, 0x1A)  # deep navy-black
CARD     = RGBColor(0x18, 0x1F, 0x28)  # card surface
CODE_BG  = RGBColor(0x20, 0x28, 0x32)  # code block
ACCENT   = RGBColor(0xF2, 0x5C, 0x3B)  # warm coral-red
GREEN    = RGBColor(0x2E, 0xCC, 0x71)  # success green
BLUE     = RGBColor(0x54, 0x9E, 0xFF)  # info blue
AMBER    = RGBColor(0xF7, 0x9E, 0x1B)  # highlight amber
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_1   = RGBColor(0xE8, 0xEC, 0xF0)  # primary text
GRAY_2   = RGBColor(0xA0, 0xAC, 0xBA)  # secondary text
GRAY_3   = RGBColor(0x6B, 0x78, 0x8A)  # tertiary text
BORDER   = RGBColor(0x2A, 0x33, 0x3F)  # subtle border

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ══════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════

def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_accent(slide, l, t, w, h=Inches(0.04), color=ACCENT):
    return add_rect(slide, l, t, w, h, color)


def add_text(slide, l, t, w, h, text, size=15, color=GRAY_1, bold=False,
             align=PP_ALIGN.LEFT, font='Microsoft YaHei', spacing=1.25,
             anchor=MSO_ANCHOR.TOP):
    """Multi-line text box. Returns shape."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = None
    tf = box.text_frame
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = Pt(size * spacing)
    return box


def add_bullets(slide, l, t, w, h, items, size=14, color=GRAY_1,
                bullet='▸', gap=Pt(4)):
    """Bullet list text box."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    tf = box.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'{bullet} {item}'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = gap
        p.line_spacing = Pt(size * 1.35)
    return box


def add_code(slide, l, t, w, h, text, size=13):
    """Styled code block."""
    add_rect(slide, l, t, w, h, CODE_BG, radius=False)
    box = slide.shapes.add_textbox(l + Inches(0.20), t + Inches(0.10),
                                   w - Inches(0.40), h - Inches(0.20))
    box.text_frame.word_wrap = True
    tf = box.text_frame
    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = GREEN
        p.font.name = 'Cascadia Code'
        p.line_spacing = Pt(size * 1.3)
    return box


def add_circle(slide, l, t, num, d=0.48):
    """Numbered circle indicator."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT
    c.line.fill.background()
    p = c.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(int(d * 28))
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return c


def add_image(slide, fname, l, t, w, h):
    """Place image scaled within box, centered. Returns shape or None."""
    path = os.path.join(IMG_DIR, fname)
    if not os.path.exists(path):
        add_text(slide, l, t + h / 3, w, Inches(0.35),
                 f'[Image: {fname}]', size=10, color=GRAY_3,
                 align=PP_ALIGN.CENTER)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    cx = l + (w - nw) // 2
    cy = t + (h - nh) // 2
    return slide.shapes.add_picture(path, cx, cy, nw, nh)


# ── Reusable header strip ─────────────────────────────────────
def slide_header(slide, step_num=None, title='', subtitle=''):
    """Draw consistent header: accent bar + optional step circle + title."""
    add_accent(slide, MARGIN_L, Inches(0.45), Inches(0.06), Inches(0.40))
    lx = MARGIN_L + Inches(0.20)
    if step_num:
        add_circle(slide, MARGIN_L, Inches(0.30), step_num, 0.42)
        lx = MARGIN_L + Inches(0.70)
    add_text(slide, lx, Inches(0.28), Inches(9.0), Inches(0.55),
             title, size=28, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, lx, Inches(0.82), Inches(9.0), Inches(0.30),
                 subtitle, size=12, color=GRAY_2)


# ── Two-column image layout helper ─────────────────────────────
def two_images(slide, y, fname1, fname2, w_each=Inches(5.6), h=Inches(3.8),
               gap=Inches(0.25)):
    """Place two images side-by-side centered on slide."""
    total_w = w_each * 2 + gap
    x1 = (SLIDE_W - total_w) // 2
    x2 = x1 + w_each + gap
    # Cards behind images
    add_rect(slide, x1 - Inches(0.06), y - Inches(0.06),
             w_each + Inches(0.12), h + Inches(0.12), CARD)
    add_rect(slide, x2 - Inches(0.06), y - Inches(0.06),
             w_each + Inches(0.12), h + Inches(0.12), CARD)
    add_image(slide, fname1, x1, y, w_each, h)
    add_image(slide, fname2, x2, y, w_each, h)


# ── Single image + text layout ─────────────────────────────────
def image_right(slide, y, fname, img_w=Inches(5.2), img_h=Inches(4.6)):
    """Place an image card on the right side."""
    ix = SLIDE_W - MARGIN_R - img_w
    add_rect(slide, ix - Inches(0.08), y - Inches(0.08),
             img_w + Inches(0.16), img_h + Inches(0.16), CARD)
    add_image(slide, fname, ix, y, img_w, img_h)
    return ix  # left edge of image area, useful for text wrapping


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 · Cover
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_accent(s, Inches(0), Inches(2.55), SLIDE_W, color=ACCENT)
add_accent(s, Inches(0), Inches(5.25), SLIDE_W, color=ACCENT)
add_text(s, Inches(1.5), Inches(2.85), Inches(10.3), Inches(1.10),
         'Claude Code 安装教程', size=52, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.95), Inches(10.3), Inches(0.60),
         '从零开始 · 轻松上手 AI 编程助手', size=22, color=GRAY_2,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.55), Inches(10.3), Inches(0.45),
         'Node.js  |  Git  |  Claude Code CLI  |  CC-Switch  |  DeepSeek  |  VS Code',
         size=14, color=GRAY_3, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.40),
         '搭载 DeepSeek 国产大模型 · 无需科学上网 · 成本仅为 Claude 官方的 2%',
         size=13, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 · Overview
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, title='安装步骤总览',
             subtitle='6 步完成 Claude Code 安装与 DeepSeek 接入')

steps = [
    ('01', 'Node.js', '运行环境，Claude Code 的基石'),
    ('02', 'Git', '版本控制工具，必须组件'),
    ('03', 'Claude Code CLI', 'npm 全局安装 + 配置文件修改'),
    ('04', 'CC-Switch', '开源模型切换工具，告别登录限制'),
    ('05', 'DeepSeek API', '国产大模型，成本仅为 Claude 的 2%'),
    ('06', 'VS Code 集成', '可视化编辑器，效率翻倍'),
]
row_h = Inches(0.72)
start_y = Inches(1.65)
left_x = MARGIN_L
list_w = Inches(7.3)

for i, (num, title, desc) in enumerate(steps):
    y = start_y + row_h * i
    add_rect(s, left_x, y, list_w, row_h - Inches(0.06), CARD)
    add_circle(s, left_x + Inches(0.12), y + Inches(0.10), num, 0.44)
    add_text(s, left_x + Inches(0.75), y + Inches(0.06),
             Inches(2.2), Inches(0.30), title, size=16, color=WHITE, bold=True)
    add_text(s, left_x + Inches(0.75), y + Inches(0.36),
             Inches(5.5), Inches(0.26), desc, size=12, color=GRAY_2)

# Right panel
rx = left_x + list_w + Inches(0.30)
rw = SLIDE_W - MARGIN_R - rx
add_rect(s, rx, start_y, rw, row_h * 6 - Inches(0.06), CARD)
add_text(s, rx + Inches(0.25), start_y + Inches(0.20), rw - Inches(0.50),
         Inches(0.35), '📌 为什么用 CC-Switch？', size=15, color=AMBER, bold=True)
add_bullets(s, rx + Inches(0.25), start_y + Inches(0.72),
            rw - Inches(0.50), Inches(3.4), [
                'Claude 是国外产品，需科学上网',
                '账号经常被封禁，不稳定',
                'API 价格昂贵 ($3~15 / 百万 tokens)',
                '',
                '✅ CC-Switch 一键切换模型',
                '✅ 接入国产 DeepSeek',
                '✅ 无需翻墙，成本降低 98%',
                '✅ 中文理解能力更出色',
            ], size=11, color=GRAY_2, bullet='•')


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 · Node.js
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='1', title='安装 Node.js —— 运行环境',
             subtitle='Claude Code 依赖 Node.js 18+，推荐安装 LTS 版本')

img_x = image_right(s, Inches(1.45), 'image.png', Inches(5.0), Inches(4.8))
text_x = MARGIN_L
text_w = img_x - text_x - Inches(0.40)
y0 = Inches(1.50)

add_text(s, text_x, y0, text_w, Inches(0.35),
         '📥 下载安装', size=20, color=AMBER, bold=True)
add_bullets(s, text_x, y0 + Inches(0.45), text_w, Inches(0.65), [
    '官网：https://nodejs.org/zh-cn/download',
    '推荐选择 LTS（长期支持）版本',
], size=13, color=GRAY_1)
add_code(s, text_x, y0 + Inches(1.20), text_w, Inches(0.40),
         'https://nodejs.org/zh-cn/download', size=12)

add_text(s, text_x, y0 + Inches(1.90), text_w, Inches(0.35),
         '✅ 验证安装', size=20, color=GREEN, bold=True)
add_bullets(s, text_x, y0 + Inches(2.35), text_w, Inches(0.45), [
    'Win + R → 输入 cmd → 回车',
], size=13, color=GRAY_1)
add_code(s, text_x, y0 + Inches(2.90), Inches(3.0), Inches(0.40),
         'node -v', size=13)
add_text(s, text_x, y0 + Inches(3.50), text_w, Inches(0.60),
         '出现版本号（如 v20.x.x）表示安装成功 ✓\n'
         '如果未出现，请重新安装或检查环境变量',
         size=13, color=GREEN, bold=False)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 · Git
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='2', title='安装 Git —— 版本控制工具',
             subtitle='Claude Code 的必要组件，不安装会直接报错')

img_x = image_right(s, Inches(1.45), 'image-2.png', Inches(5.0), Inches(4.8))
text_x = MARGIN_L
text_w = img_x - text_x - Inches(0.40)
y0 = Inches(1.50)

add_text(s, text_x, y0, text_w, Inches(0.35),
         '📥 下载与安装', size=20, color=AMBER, bold=True)
add_bullets(s, text_x, y0 + Inches(0.45), text_w, Inches(1.20), [
    '官网：https://git-scm.com/downloads',
    '下载 Windows 版本安装包',
    '双击运行，保持默认选项一路 Next',
    '安装后 cmd 中输入 git --version 验证',
], size=13, color=GRAY_1)

add_text(s, text_x, y0 + Inches(2.25), text_w, Inches(0.35),
         '⚠️ 重要提醒', size=18, color=ACCENT, bold=True)
add_bullets(s, text_x, y0 + Inches(2.70), text_w, Inches(1.00), [
    '不装 Git，Claude Code 直接报错退出',
    'Git 是免费开源软件，安全可靠',
    '如已安装仍报错，检查 PATH 环境变量',
], size=13, color=GRAY_1)
add_text(s, text_x, y0 + Inches(3.90), text_w, Inches(0.50),
         '✅ Git 安装成功后，Claude Code 已经可以运行！\n'
         '   但还需要修改配置文件跳过登录验证……',
         size=13, color=GREEN, bold=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 · Claude Code CLI
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='3', title='安装 Claude Code CLI',
             subtitle='npm 全局安装 + 配置文件修改，两步完成')

img_x = image_right(s, Inches(1.45), 'image-1.png', Inches(5.0), Inches(4.8))
text_x = MARGIN_L
text_w = img_x - text_x - Inches(0.40)
y0 = Inches(1.50)

add_text(s, text_x, y0, text_w, Inches(0.30),
         '📦 安装命令', size=20, color=AMBER, bold=True)
add_bullets(s, text_x, y0 + Inches(0.40), text_w, Inches(0.40), [
    '使用 PowerShell（管理员身份）或 cmd 执行',
], size=12, color=GRAY_1)
add_code(s, text_x, y0 + Inches(0.90), text_w + Inches(0.4), Inches(0.40),
         'npm install -g @anthropic-ai/claude-code', size=13)

add_text(s, text_x, y0 + Inches(1.55), text_w, Inches(0.30),
         '💡 加速技巧（安装太慢时使用）', size=16, color=BLUE, bold=True)
add_code(s, text_x, y0 + Inches(1.95), text_w + Inches(0.4), Inches(0.35),
         'npm config set registry https://registry.npmmirror.com', size=12)

add_text(s, text_x, y0 + Inches(2.60), text_w, Inches(0.30),
         '🔧 修改配置文件', size=20, color=ACCENT, bold=True)
add_bullets(s, text_x, y0 + Inches(3.00), text_w, Inches(0.55), [
    '找到 C:\\Users\\用户名\\.claude\\claude.json',
    '添加 "hasCompletedOnboarding": true',
], size=12, color=GRAY_1)
add_code(s, text_x, y0 + Inches(3.65), Inches(3.2), Inches(0.60),
         '{\n  "hasCompletedOnboarding": true\n}', size=13)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 · First run & config
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='3', title='首次运行 —— 登录提示',
             subtitle='当前 Claude Code 能运行，但提示登录 Claude 账号')

y_img = Inches(1.45)
two_images(s, y_img, 'image-5.png', 'image-6.png',
           w_each=Inches(5.4), h=Inches(3.4), gap=Inches(0.30))

y_bot = Inches(5.20)
add_text(s, MARGIN_L, y_bot, content_w := SLIDE_W - MARGIN_L * 2,
         Inches(0.35), '⚠️ 当前状态', size=18, color=ACCENT, bold=True)
add_text(s, MARGIN_L, y_bot + Inches(0.40), Inches(5.5), Inches(0.80),
         '已完成 claude.json 配置，可以运行\n'
         '但会提示需要登录 Claude 账号',
         size=13, color=GRAY_2)

add_text(s, Inches(7.20), y_bot + Inches(0.40), Inches(5.5), Inches(1.30),
         '💡 问题分析：\n\n'
         '   Claude 是国外大模型，需要：\n'
         '   · 科学上网才能访问\n'
         '   · 付费购买 token（价格昂贵）\n'
         '   · 账号经常被封\n\n'
         '   👉 解决方案：接入国产 DeepSeek！',
         size=12, color=GRAY_2)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 · CC-Switch Introduction
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='4', title='CC-Switch —— 一键切换大模型',
             subtitle='开源工具，图形界面操作，支持 DeepSeek 等多款国产模型')

y_img = Inches(1.45)
two_images(s, y_img, 'image-7.png', 'image-8.png',
           w_each=Inches(5.4), h=Inches(3.4), gap=Inches(0.30))

y_bot = Inches(5.20)
add_text(s, MARGIN_L, y_bot, SLIDE_W - MARGIN_L * 2, Inches(0.30),
         '🔑 CC-Switch 核心优势', size=18, color=AMBER, bold=True)
add_bullets(s, MARGIN_L, y_bot + Inches(0.40),
            SLIDE_W - MARGIN_L * 2, Inches(1.60), [
                '开源免费，GitHub 托管：https://github.com/farion1231/cc-switch',
                '一键切换 Claude Code 底层大模型 API，无需修改任何代码',
                '支持 DeepSeek、通义千问等多个国产模型，图形界面填写 API Key 即可',
                '下载 Releases 中的 .exe 安装包，双击运行安装，无需额外配置',
            ], size=13, color=GRAY_1)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 · DeepSeek API Key
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='5', title='获取 DeepSeek API 密钥',
             subtitle='在 DeepSeek 开放平台注册账号，创建 API Key')

y_img = Inches(1.45)
two_images(s, y_img, 'image-9.png', 'image-10.png',
           w_each=Inches(5.4), h=Inches(2.8), gap=Inches(0.30))

y_bot = Inches(4.60)
add_text(s, MARGIN_L, y_bot, SLIDE_W - MARGIN_L * 2, Inches(0.30),
         '📋 操作步骤', size=18, color=AMBER, bold=True)
add_bullets(s, MARGIN_L, y_bot + Inches(0.40),
            Inches(6.5), Inches(1.60), [
                '① 打开 platform.deepseek.com → 手机号注册',
                '② 登录后进入「API Keys」管理页面',
                '③ 点击「创建 API Key」→ 起个名字（如 claude-code）',
                '④ ⚠️ 立即复制密钥（sk-xxxx），仅显示一次！',
            ], size=13, color=GRAY_1)

add_rect(s, Inches(7.5), y_bot + Inches(0.40), Inches(5.3), Inches(1.60), CARD)
add_text(s, Inches(7.7), y_bot + Inches(0.55), Inches(4.9), Inches(1.30),
         '🔐 安全提醒\n\n'
         '· API Key 关闭页面后无法再次查看\n'
         '· 请复制后保存到安全的地方\n'
         '· 不要将 Key 分享给他人\n'
         '· 泄露后可在平台重新生成',
         size=12, color=GRAY_2)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 · CC-Switch Configuration
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, title='配置 CC-Switch 接入 DeepSeek',
             subtitle='添加模型 → 选择 DeepSeek → 粘贴 API Key，完成！')

y_img = Inches(1.45)
two_images(s, y_img, 'image-11.png', 'image-12.png',
           w_each=Inches(5.4), h=Inches(3.0), gap=Inches(0.30))

y_bot = Inches(4.80)
add_text(s, MARGIN_L, y_bot, SLIDE_W - MARGIN_L * 2, Inches(0.30),
         '🎉 配置完成！重新打开 cmd，输入 claude 不再提示登录',
         size=16, color=GREEN, bold=True)
add_bullets(s, MARGIN_L, y_bot + Inches(0.40),
            SLIDE_W - MARGIN_L * 2, Inches(1.50), [
                '① 打开 CC-Switch → 进入主界面',
                '② 点击「添加模型」→ 在模型列表中选择「DeepSeek」',
                '③ 在 API Key 输入框中粘贴刚才复制的密钥 → 保存',
                '④ 重新打开 cmd → 输入 claude → 直接进入对话界面！',
            ], size=13, color=GRAY_1)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 · DeepSeek Pricing & Topup
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='5', title='DeepSeek 定价与充值',
             subtitle='实名认证 + 充值 10 元即可，成本仅为 Claude 官方的 2%')

left_x = MARGIN_L
y0 = Inches(1.50)
card_w = Inches(3.6)
card_h = Inches(1.80)

# ── DeepSeek pricing card ──
add_rect(s, left_x, y0, card_w, card_h, CARD)
add_accent(s, left_x, y0, card_w, color=GREEN)
add_text(s, left_x + Inches(0.25), y0 + Inches(0.15), Inches(3.0), Inches(0.30),
         '🟢 DeepSeek-V3', size=18, color=GREEN, bold=True)
add_text(s, left_x + Inches(0.25), y0 + Inches(0.60), Inches(3.0), Inches(1.10),
         '输入：¥1.0  / 百万 tokens\n'
         '输出：¥2.0  / 百万 tokens\n\n'
         '≈ 充值 1 元 = 50 万字对话\n'
         '≈ 充值 10 元 = 500 万字对话',
         size=13, color=GRAY_1)

# ── Claude pricing card ──
cx = left_x + card_w + Inches(0.40)
add_rect(s, cx, y0, card_w, card_h, CARD)
add_accent(s, cx, y0, card_w, color=RGBColor(0xEF, 0x44, 0x44))
add_text(s, cx + Inches(0.25), y0 + Inches(0.15), Inches(3.0), Inches(0.30),
         '🔴 Claude Sonnet 4', size=18, color=RGBColor(0xEF, 0x44, 0x44), bold=True)
add_text(s, cx + Inches(0.25), y0 + Inches(0.60), Inches(3.0), Inches(1.10),
         '输入：$3 ≈ ¥21  / 百万 tokens\n'
         '输出：$15 ≈ ¥105 / 百万 tokens\n\n'
         '≈ 价格是 DeepSeek 的 50 倍+\n'
         '≈ 还需科学上网 + 国外信用卡',
         size=13, color=GRAY_1)

# ── Cost summary bar ──
add_text(s, left_x, y0 + card_h + Inches(0.20), Inches(7.5), Inches(0.30),
         '👉 结论：DeepSeek 成本仅为 Claude 的 2%，且中文理解能力更强！',
         size=14, color=GREEN, bold=True)

# ── Right side: 3 images ──
ix = Inches(8.5)
img_w = Inches(1.45)
img_h = Inches(2.0)
add_text(s, ix, y0, Inches(4.5), Inches(0.25),
         '充值流程截图', size=12, color=GRAY_3, align=PP_ALIGN.CENTER)
for j, fname in enumerate(['image-13.png', 'image-14.png', 'image-15.png']):
    px = ix + (img_w + Inches(0.12)) * j
    add_rect(s, px - Inches(0.04), y0 + Inches(0.35) - Inches(0.04),
             img_w + Inches(0.08), img_h + Inches(0.08), CARD)
    add_image(s, fname, px, y0 + Inches(0.35), img_w, img_h)

# ── Bottom: steps ──
y1 = y0 + card_h + Inches(0.75)
add_text(s, left_x, y1, Inches(4.0), Inches(0.28),
         '📝 充值步骤', size=18, color=BLUE, bold=True)
add_bullets(s, left_x, y1 + Inches(0.38), Inches(7.8), Inches(1.30), [
    '① 在 DeepSeek 平台完成实名认证（按平台指引操作）',
    '② 进入充值页面，首次建议充值 10~20 元',
    '③ 支持微信 / 支付宝扫码支付，即时到账',
    '④ 充值后即可无限制调用 API，按量计费，用多少扣多少',
], size=13, color=GRAY_1)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 · VS Code
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, step_num='6', title='VS Code 集成 Claude Code',
             subtitle='在可视化编辑器中运行，告别纯命令行，效率翻倍')

img_x = image_right(s, Inches(1.45), 'image-16.png', Inches(5.0), Inches(4.6))
text_x = MARGIN_L
text_w = img_x - text_x - Inches(0.35)
y0 = Inches(1.45)

add_text(s, text_x, y0, text_w, Inches(0.30),
         '🌟 为什么选择 VS Code？', size=18, color=AMBER, bold=True)
add_bullets(s, text_x, y0 + Inches(0.40), text_w, Inches(1.00), [
    '微软出品的免费编辑器，轻量且功能强大',
    '可视化操作界面，告别纯命令行',
    '可直观查看代码变更、对比差异',
    '丰富插件生态，无限扩展可能',
], size=12, color=GRAY_1)

add_text(s, text_x, y0 + Inches(1.70), text_w, Inches(0.30),
         '🔌 安装 Claude Code 插件', size=18, color=BLUE, bold=True)
add_bullets(s, text_x, y0 + Inches(2.10), text_w, Inches(1.10), [
    'Ctrl+Shift+X → 搜索 "Claude Code"',
    '找到 Anthropic 官方插件 → Install',
    'Ctrl+` 打开终端 → 输入 claude',
], size=12, color=GRAY_1)

add_text(s, text_x, y0 + Inches(3.50), text_w, Inches(0.30),
         '⚠️ PowerShell 执行策略报错？', size=16, color=ACCENT, bold=True)
add_code(s, text_x, y0 + Inches(3.90), text_w + Inches(0.40), Inches(0.45),
         'Set-ExecutionPolicy -ExecutionPolicy RemoteSigned '
         '-Scope CurrentUser', size=11)
add_text(s, text_x, y0 + Inches(4.50), text_w, Inches(0.40),
         '在 PowerShell 中执行上述命令，输入 Y 确认即可\n'
         'RemoteSigned：允许本地脚本运行，保留安全防护',
         size=11, color=GRAY_2)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 · PowerShell Policy Fix (dedicated)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, title='解决 PowerShell 执行策略报错',
             subtitle='claude.ps1 无法运行？这是 Windows 安全机制，一行命令解决')

text_x = MARGIN_L
text_w = Inches(6.8)
y0 = Inches(1.50)

add_text(s, text_x, y0, text_w, Inches(0.30),
         '🔍 问题分析', size=20, color=AMBER, bold=True)
add_bullets(s, text_x, y0 + Inches(0.42), text_w, Inches(1.10), [
    'Windows PowerShell 默认执行策略为 Restricted',
    '该策略禁止所有脚本运行，包括 claude.ps1',
    '这是 Windows 安全功能，防止恶意脚本执行',
    '修改仅影响当前用户，不影响系统安全',
], size=13, color=GRAY_1)

add_text(s, text_x, y0 + Inches(1.85), text_w, Inches(0.30),
         '✅ 解决方案（推荐）', size=20, color=GREEN, bold=True)
add_bullets(s, text_x, y0 + Inches(2.27), text_w, Inches(0.55), [
    '以普通身份打开 PowerShell（无需管理员权限）',
    '执行以下命令（仅修改当前用户策略）',
], size=12, color=GRAY_1)
add_code(s, text_x, y0 + Inches(2.95), text_w, Inches(0.50),
         'Set-ExecutionPolicy -ExecutionPolicy RemoteSigned '
         '-Scope CurrentUser', size=13)

add_text(s, text_x, y0 + Inches(3.70), text_w, Inches(0.70),
         'RemoteSigned 策略说明：\n'
         '   ✅ 允许运行本地创建的脚本\n'
         '   ⚠️ 从网络下载的脚本仍需有效签名\n'
         '   🔒 既解决问题，又保留基本安全防护\n\n'
         '系统询问 "是否要更改执行策略?" → 输入 Y 回车确认\n'
         '关闭当前终端 → 重新打开 → 再次运行 claude 即可',
         size=12, color=GRAY_2)

# Right image
img_x = Inches(7.9)
img_w = Inches(5.0)
add_rect(s, img_x - Inches(0.08), Inches(1.45) - Inches(0.08),
         img_w + Inches(0.16), Inches(5.0) + Inches(0.16), CARD)
add_image(s, 'image-17.png', img_x, Inches(1.45), img_w, Inches(5.0))


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 · FAQ / Troubleshooting
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_header(s, title='常见问题排查',
             subtitle='安装过程中遇到的典型问题及解决方案')

faqs = [
    ('Q1: npm install 报错',
     '→ 切换国内镜像：npm config set registry\n'
     '  https://registry.npmmirror.com\n'
     '→ 或用管理员身份运行 PowerShell'),
    ('Q2: 提示找不到 claude 命令',
     '→ 关闭并重新打开 cmd 窗口\n'
     '→ 确认 npm 全局安装成功\n'
     '→ 检查 npm 全局路径是否在 PATH 中'),
    ('Q3: 提示找不到 Git',
     '→ 确认已安装：git --version\n'
     '→ 重装 Git 时勾选「Add Git to PATH」\n'
     '→ 或手动添加 Git 到系统环境变量'),
    ('Q4: CC-Switch 配置后仍提示登录',
     '→ 检查 API Key 是否有多余空格\n'
     '→ 确认 DeepSeek 已实名认证并充值\n'
     '→ 重启 CC-Switch 和 cmd 再试'),
    ('Q5: DeepSeek API 返回错误',
     '→ 检查账户余额是否充足\n'
     '→ 确认 API Key 未被删除/过期\n'
     '→ 实名认证是否已审核通过'),
    ('Q6: VS Code 终端无法运行 claude',
     '→ 按前一页修改 PowerShell 执行策略\n'
     '→ 关闭 VS Code 重新打开\n'
     '→ 确认外部 cmd 中可正常运行'),
]

n_cols = 3
card_w = Inches(3.75)
card_h = Inches(1.70)
gap_x = Inches(0.20)
gap_y = Inches(0.15)
total_w = card_w * n_cols + gap_x * (n_cols - 1)
start_x = (SLIDE_W - total_w) // 2
start_y = Inches(1.55)

for i, (title, desc) in enumerate(faqs):
    col = i % n_cols
    row = i // n_cols
    cx = start_x + (card_w + gap_x) * col
    cy = start_y + (card_h + gap_y) * row
    add_rect(s, cx, cy, card_w, card_h, CARD)
    add_accent(s, cx, cy, card_w, color=BLUE)
    add_text(s, cx + Inches(0.15), cy + Inches(0.12),
             Inches(3.4), Inches(0.28),
             title, size=15, color=AMBER, bold=True)
    add_text(s, cx + Inches(0.15), cy + Inches(0.48),
             Inches(3.4), Inches(1.15),
             desc, size=11, color=GRAY_2)

# Bottom hint
add_text(s, MARGIN_L, Inches(6.35), SLIDE_W - MARGIN_L * 2, Inches(0.30),
         '💡 更多帮助请查看完整教程：Claude_Code_安装教程.md',
         size=11, color=GRAY_3, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 · Summary
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_accent(s, Inches(0), Inches(1.80), SLIDE_W, color=ACCENT)

add_text(s, Inches(1.5), Inches(0.55), Inches(10.3), Inches(0.75),
         '🎯 安装完成！', size=42, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(1.25), Inches(10.3), Inches(0.35),
         '在终端输入 claude 即可开启 AI 编程之旅',
         size=15, color=GRAY_2, align=PP_ALIGN.CENTER)

cards = [
    ('Node.js', '运行环境', 'nodejs.org', ACCENT),
    ('Git', '版本控制', 'git-scm.com', AMBER),
    ('Claude Code', 'AI 编程 CLI', 'npm install -g', BLUE),
    ('CC-Switch', '模型切换', 'GitHub', GREEN),
    ('DeepSeek', '国产大模型', 'platform.deepseek', RGBColor(0x4D, 0x6B, 0xFF)),
    ('VS Code', '可视化编辑', '插件市场', ACCENT),
]
card_w = Inches(1.85)
card_h = Inches(3.00)
n = len(cards)
gap = Inches(0.18)
total_w = card_w * n + gap * (n - 1)
start_x = (SLIDE_W - total_w) // 2
cy = Inches(2.40)

for i, (name, desc, link, clr) in enumerate(cards):
    cx = start_x + (card_w + gap) * i
    add_rect(s, cx, cy, card_w, card_h, CARD)
    add_accent(s, cx, cy, card_w, color=clr)
    add_text(s, cx + Inches(0.10), cy + Inches(0.25),
             Inches(1.65), Inches(0.35),
             name, size=16, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx + Inches(0.10), cy + Inches(1.00),
             Inches(1.65), Inches(0.35),
             desc, size=11, color=GRAY_2, align=PP_ALIGN.CENTER)
    add_text(s, cx + Inches(0.10), cy + Inches(1.75),
             Inches(1.65), Inches(0.35),
             link, size=9, color=GRAY_3, align=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.80),
         'Claude Code + DeepSeek = 高效 · 低成本 · 无限制\n'
         '无需科学上网 · 无需国外账号 · 成本仅为 Claude 官方的 2%',
         size=17, color=GRAY_1, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'PPT saved: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
